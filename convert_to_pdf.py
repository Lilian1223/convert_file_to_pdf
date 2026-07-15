"""
拖放轉 PDF 小工具（免安裝 LibreOffice / Office 版）
--------------------------------------------------
把檔案拖曳到這支程式（或它的捷徑）的圖示上，
會自動把檔案轉成 PDF 並存到桌面。

這個版本完全用 Python 套件轉換，不需要另外安裝
LibreOffice 或 Microsoft Office，但代價是：
  - 排版是「重新產生」的簡化版，不會 100% 還原原本的樣式、字型、圖片位置
  - 只支援新版 Office 格式：docx / xlsx / pptx（舊版 .doc/.xls/.ppt 不支援）
  - 圖片轉 PDF 是完整品質（img2pdf 為無損轉換）

支援格式：
  docx, txt, xlsx, csv, pptx, jpg/jpeg/png/bmp/gif/tiff, pdf(直接複製)

安裝需求（只需執行一次）：
  pip install python-docx openpyxl python-pptx img2pdf reportlab pillow
"""

import sys
import os
import shutil
import traceback
import ctypes
import ctypes.wintypes
from xml.sax.saxutils import escape as xml_escape

# 用 pythonw.exe 執行時沒有主控台，sys.stdout/stderr 會是 None，
# 若有任何底層程式呼叫 print() 會直接出錯，這裡先接到黑洞避免整支程式壞掉
if sys.stdout is None:
    sys.stdout = open(os.devnull, "w")
if sys.stderr is None:
    sys.stderr = open(os.devnull, "w")


def get_desktop_path():
    CSIDL_DESKTOPDIRECTORY = 0x0010
    buf = ctypes.create_unicode_buffer(ctypes.wintypes.MAX_PATH)
    ctypes.windll.shell32.SHGetFolderPathW(None, CSIDL_DESKTOPDIRECTORY, None, 0, buf)
    return buf.value


def unique_path(path):
    if not os.path.exists(path):
        return path
    base, ext = os.path.splitext(path)
    i = 1
    while os.path.exists(f"{base}({i}){ext}"):
        i += 1
    return f"{base}({i}){ext}"


def get_cjk_styles():
    """建立支援中文（繁體）顯示的段落樣式"""
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.enums import TA_LEFT

    font_name = "MSung-Light"
    if font_name not in pdfmetrics.getRegisteredFontNames():
        pdfmetrics.registerFont(UnicodeCIDFont(font_name))

    normal = ParagraphStyle("CJKNormal", fontName=font_name, fontSize=10, leading=15, alignment=TA_LEFT)
    heading = ParagraphStyle("CJKHeading", fontName=font_name, fontSize=14, leading=20, spaceAfter=8)
    mono = ParagraphStyle("CJKMono", fontName=font_name, fontSize=9, leading=13)
    return normal, heading, mono, font_name


def to_cell(text, style):
    from reportlab.platypus import Paragraph
    text = "" if text is None else str(text)
    return Paragraph(xml_escape(text).replace("\n", "<br/>"), style)


def docx_to_pdf(src, dest):
    from docx import Document
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib import colors

    normal, heading, mono, font_name = get_cjk_styles()
    doc = Document(src)
    story = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            story.append(Spacer(1, 0.25 * cm))
            continue
        style_name = para.style.name if para.style else ""
        style = heading if "Heading" in style_name or "Title" in style_name else normal
        story.append(Paragraph(xml_escape(text), style))
        story.append(Spacer(1, 0.1 * cm))

    for table in doc.tables:
        data = [[to_cell(cell.text, normal) for cell in row.cells] for row in table.rows]
        t = Table(data)
        t.setStyle(TableStyle([("GRID", (0, 0), (-1, -1), 0.5, colors.grey)]))
        story.append(t)
        story.append(Spacer(1, 0.3 * cm))

    if not story:
        story.append(Paragraph("（空白文件）", normal))

    SimpleDocTemplate(dest, pagesize=A4).build(story)


def txt_to_pdf(src, dest):
    from reportlab.platypus import SimpleDocTemplate, Paragraph
    from reportlab.lib.pagesizes import A4

    normal, heading, mono, font_name = get_cjk_styles()
    with open(src, encoding="utf-8", errors="replace") as f:
        content = f.read()
    story = [Paragraph(xml_escape(content).replace("\n", "<br/>"), mono)]
    SimpleDocTemplate(dest, pagesize=A4).build(story)


def xlsx_to_pdf(src, dest):
    from openpyxl import load_workbook
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Table, TableStyle, PageBreak, Spacer
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib import colors
    from reportlab.lib.units import cm

    normal, heading, mono, font_name = get_cjk_styles()
    wb = load_workbook(src, data_only=True)
    story = []

    for sheet in wb.worksheets:
        story.append(Paragraph(xml_escape(sheet.title), heading))
        rows = list(sheet.iter_rows(values_only=True))
        if rows:
            data = [[to_cell(c, normal) for c in row] for row in rows]
            t = Table(data)
            t.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]))
            story.append(t)
        else:
            story.append(Paragraph("（空白工作表）", normal))
        story.append(PageBreak())

    if story and isinstance(story[-1], type(story[0])) is False:
        pass
    if story:
        story.pop()  # 移除最後多餘的換頁

    SimpleDocTemplate(dest, pagesize=landscape(A4)).build(story or [Paragraph("（空白活頁簿）", normal)])


def csv_to_pdf(src, dest):
    import csv
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib import colors

    normal, heading, mono, font_name = get_cjk_styles()
    with open(src, encoding="utf-8-sig", errors="replace", newline="") as f:
        rows = list(csv.reader(f))
    data = [[to_cell(c, normal) for c in row] for row in rows] or [[to_cell("", normal)]]
    t = Table(data)
    t.setStyle(TableStyle([("GRID", (0, 0), (-1, -1), 0.5, colors.grey)]))
    SimpleDocTemplate(dest, pagesize=landscape(A4)).build([t])


def pptx_to_pdf(src, dest):
    from pptx import Presentation
    from reportlab.platypus import SimpleDocTemplate, Paragraph, PageBreak
    from reportlab.lib.pagesizes import A4

    normal, heading, mono, font_name = get_cjk_styles()
    prs = Presentation(src)
    story = []

    for i, slide in enumerate(prs.slides, 1):
        story.append(Paragraph(f"投影片 {i}", heading))
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        story.append(Paragraph(xml_escape(text), normal))
        story.append(PageBreak())

    if story:
        story.pop()

    SimpleDocTemplate(dest, pagesize=A4).build(story or [Paragraph("（空白簡報）", normal)])


def image_to_pdf(src, dest):
    import img2pdf
    with open(dest, "wb") as f:
        f.write(img2pdf.convert(src))


CONVERTERS = {
    ".docx": docx_to_pdf,
    ".txt": txt_to_pdf,
    ".xlsx": xlsx_to_pdf,
    ".csv": csv_to_pdf,
    ".pptx": pptx_to_pdf,
    ".jpg": image_to_pdf,
    ".jpeg": image_to_pdf,
    ".png": image_to_pdf,
    ".bmp": image_to_pdf,
    ".gif": image_to_pdf,
    ".tiff": image_to_pdf,
}

LEGACY_UNSUPPORTED = {".doc", ".xls", ".ppt"}


def convert_to_pdf(filepath, desktop):
    filepath = os.path.abspath(filepath)
    if not os.path.isfile(filepath):
        return False, f"找不到檔案：{filepath}"

    ext = os.path.splitext(filepath)[1].lower()
    dest = unique_path(os.path.join(desktop, os.path.splitext(os.path.basename(filepath))[0] + ".pdf"))

    if ext == ".pdf":
        dest = unique_path(os.path.join(desktop, os.path.basename(filepath)))
        shutil.copy2(filepath, dest)
        return True, dest

    if ext in LEGACY_UNSUPPORTED:
        return False, f"不支援舊版格式 {ext}，請先在 Word/Excel/PowerPoint 另存為新格式（{ext}x）再拖曳"

    converter = CONVERTERS.get(ext)
    if converter is None:
        return False, f"不支援的檔案類型：{ext}"

    converter(filepath, dest)
    return True, dest


def show_message_box(title, text, icon=0x40):
    """跳出 Windows 訊息視窗（0x40=資訊, 0x10=錯誤, 0x30=警告）"""
    ctypes.windll.user32.MessageBoxW(0, text, title, icon)


def show_result_notification(results):
    success = [r for r in results if r[1]]
    fail = [r for r in results if not r[1]]

    lines = []
    for src, ok, msg in results:
        name = os.path.basename(src)
        lines.append(f"[成功] {name}" if ok else f"[失敗] {name} - {msg}")

    title = f"轉換完成（成功 {len(success)}，失敗 {len(fail)}）"
    icon = 0x40 if not fail else 0x30  # 全部成功用資訊圖示，有失敗用警告圖示
    show_message_box(title, "\n".join(lines), icon)


def main():
    args = sys.argv[1:]

    # 從參數中取出模式旗標；拖放時這個旗標會固定寫在捷徑目標裡，
    # 拖上去的檔案路徑會被 Windows 自動接在後面
    notify = False
    if "--notify" in args:
        notify = True
        args = [a for a in args if a != "--notify"]
    elif "--silent" in args:
        notify = False
        args = [a for a in args if a != "--silent"]

    files = args

    if not files:
        if notify:
            show_message_box("拖放轉 PDF", "請把一個或多個檔案拖曳到這個捷徑上。", 0x30)
        return

    try:
        import docx, openpyxl, pptx, img2pdf, reportlab  # noqa: F401
    except ImportError as e:
        if notify:
            show_message_box(
                "缺少必要套件",
                "請先在命令提示字元執行：\n"
                "pip install python-docx openpyxl python-pptx img2pdf reportlab pillow\n\n"
                f"詳細錯誤：{e}",
                0x10,
            )
        return

    desktop = get_desktop_path()
    results = []
    for f in files:
        try:
            ok, msg = convert_to_pdf(f, desktop)
        except Exception as e:
            ok, msg = False, str(e)
        results.append((f, ok, msg))

    if notify:
        show_result_notification(results)


if __name__ == "__main__":
    main()
